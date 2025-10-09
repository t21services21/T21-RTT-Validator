"""
T21 UNIVERSAL DOCUMENT PROCESSOR
Extract text from ANY document type

Supported Formats:
- PDF documents
- Word documents (DOCX, DOC)
- Excel spreadsheets (XLSX, XLS)
- PowerPoint presentations (PPTX, PPT)
- Images (PNG, JPG, JPEG) with OCR
- Text files (TXT, CSV, JSON)
- HTML files
- Markdown files

Features:
- Automatic format detection
- Text extraction
- Table extraction from Excel/PDF
- OCR for images
- Preserves structure
"""

import os
import json
from typing import Dict, List, Tuple


def extract_text_from_file(file_path: str, file_type: str = None) -> Dict:
    """
    Extract text from any file type
    
    Args:
        file_path: Path to file
        file_type: Optional file type hint
        
    Returns:
        Dict with extracted text and metadata
    """
    
    # Detect file type from extension if not provided
    if not file_type:
        file_extension = os.path.splitext(file_path)[1].lower()
    else:
        file_extension = file_type.lower()
    
    result = {
        'success': False,
        'text': '',
        'file_type': file_extension,
        'error': None,
        'metadata': {}
    }
    
    try:
        # Route to appropriate handler
        if file_extension in ['.txt', '.text']:
            result = extract_from_text(file_path)
        elif file_extension in ['.pdf']:
            result = extract_from_pdf(file_path)
        elif file_extension in ['.docx', '.doc']:
            result = extract_from_word(file_path)
        elif file_extension in ['.xlsx', '.xls']:
            result = extract_from_excel(file_path)
        elif file_extension in ['.pptx', '.ppt']:
            result = extract_from_powerpoint(file_path)
        elif file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff']:
            result = extract_from_image(file_path)
        elif file_extension in ['.csv']:
            result = extract_from_csv(file_path)
        elif file_extension in ['.json']:
            result = extract_from_json(file_path)
        elif file_extension in ['.html', '.htm']:
            result = extract_from_html(file_path)
        elif file_extension in ['.md', '.markdown']:
            result = extract_from_markdown(file_path)
        else:
            result['error'] = f"Unsupported file type: {file_extension}"
            
    except Exception as e:
        result['error'] = str(e)
    
    return result


def extract_from_text(file_path: str) -> Dict:
    """Extract text from TXT file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return {
            'success': True,
            'text': text,
            'file_type': '.txt',
            'metadata': {
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_pdf(file_path: str) -> Dict:
    """Extract text from PDF file"""
    
    try:
        # Try PyPDF2 first (most common)
        try:
            import PyPDF2
            
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text() + "\n\n"
            
            return {
                'success': True,
                'text': text,
                'file_type': '.pdf',
                'metadata': {
                    'pages': num_pages,
                    'character_count': len(text),
                    'word_count': len(text.split())
                }
            }
        except ImportError:
            pass
        
        # Try pdfplumber (better for tables)
        try:
            import pdfplumber
            
            text = ""
            with pdfplumber.open(file_path) as pdf:
                num_pages = len(pdf.pages)
                
                for page in pdf.pages:
                    text += page.extract_text() + "\n\n"
            
            return {
                'success': True,
                'text': text,
                'file_type': '.pdf',
                'metadata': {
                    'pages': num_pages,
                    'character_count': len(text),
                    'word_count': len(text.split())
                }
            }
        except ImportError:
            pass
        
        # Libraries not available
        return {
            'success': False,
            'text': '',
            'error': 'PDF library not installed. Please install: pip install PyPDF2 or pip install pdfplumber',
            'instructions': 'To enable PDF support, install: pip install PyPDF2'
        }
        
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_word(file_path: str) -> Dict:
    """Extract text from Word document"""
    
    try:
        import docx
        
        doc = docx.Document(file_path)
        
        # Extract text from paragraphs
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\t"
                text += "\n"
        
        return {
            'success': True,
            'text': text,
            'file_type': '.docx',
            'metadata': {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except ImportError:
        return {
            'success': False,
            'text': '',
            'error': 'python-docx library not installed',
            'instructions': 'To enable Word support, install: pip install python-docx'
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_excel(file_path: str) -> Dict:
    """Extract text from Excel spreadsheet"""
    
    try:
        import openpyxl
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"=== SHEET: {sheet_name} ===\n\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text + "\n"
            
            text += "\n\n"
        
        return {
            'success': True,
            'text': text,
            'file_type': '.xlsx',
            'metadata': {
                'sheets': len(workbook.sheetnames),
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except ImportError:
        return {
            'success': False,
            'text': '',
            'error': 'openpyxl library not installed',
            'instructions': 'To enable Excel support, install: pip install openpyxl'
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_powerpoint(file_path: str) -> Dict:
    """Extract text from PowerPoint presentation"""
    
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        text = ""
        slide_num = 0
        
        for slide in prs.slides:
            slide_num += 1
            text += f"=== SLIDE {slide_num} ===\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            text += "\n\n"
        
        return {
            'success': True,
            'text': text,
            'file_type': '.pptx',
            'metadata': {
                'slides': len(prs.slides),
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except ImportError:
        return {
            'success': False,
            'text': '',
            'error': 'python-pptx library not installed',
            'instructions': 'To enable PowerPoint support, install: pip install python-pptx'
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_image(file_path: str) -> Dict:
    """Extract text from image using OCR"""
    
    try:
        from PIL import Image
        import pytesseract
        
        # Open image
        image = Image.open(file_path)
        
        # Perform OCR
        text = pytesseract.image_to_string(image)
        
        return {
            'success': True,
            'text': text,
            'file_type': os.path.splitext(file_path)[1],
            'metadata': {
                'image_size': image.size,
                'image_mode': image.mode,
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except ImportError:
        return {
            'success': False,
            'text': '',
            'error': 'OCR libraries not installed',
            'instructions': 'To enable image OCR, install: pip install pytesseract pillow'
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_csv(file_path: str) -> Dict:
    """Extract text from CSV file"""
    
    try:
        import csv
        
        text = ""
        with open(file_path, 'r', encoding='utf-8') as f:
            csv_reader = csv.reader(f)
            for row in csv_reader:
                text += "\t".join(row) + "\n"
        
        return {
            'success': True,
            'text': text,
            'file_type': '.csv',
            'metadata': {
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_json(file_path: str) -> Dict:
    """Extract text from JSON file"""
    
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Convert JSON to readable text
        text = json.dumps(data, indent=2)
        
        return {
            'success': True,
            'text': text,
            'file_type': '.json',
            'metadata': {
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_html(file_path: str) -> Dict:
    """Extract text from HTML file"""
    
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return {
            'success': True,
            'text': text,
            'file_type': '.html',
            'metadata': {
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except ImportError:
        return {
            'success': False,
            'text': '',
            'error': 'BeautifulSoup library not installed',
            'instructions': 'To enable HTML support, install: pip install beautifulsoup4'
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_markdown(file_path: str) -> Dict:
    """Extract text from Markdown file"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return {
            'success': True,
            'text': text,
            'file_type': '.md',
            'metadata': {
                'character_count': len(text),
                'word_count': len(text.split())
            }
        }
    except Exception as e:
        return {'success': False, 'text': '', 'error': str(e)}


def extract_from_uploaded_file(uploaded_file) -> Dict:
    """
    Extract text from Streamlit uploaded file
    
    Args:
        uploaded_file: Streamlit UploadedFile object
        
    Returns:
        Dict with extracted text and metadata
    """
    
    # Save uploaded file temporarily
    import tempfile
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=uploaded_file.name) as tmp_file:
        tmp_file.write(uploaded_file.getvalue())
        tmp_path = tmp_file.name
    
    try:
        # Extract text from temp file
        result = extract_text_from_file(tmp_path)
        result['original_filename'] = uploaded_file.name
        return result
    finally:
        # Clean up temp file
        try:
            os.unlink(tmp_path)
        except:
            pass


def get_supported_formats() -> List[str]:
    """Get list of supported file formats"""
    return [
        '.txt', '.text',
        '.pdf',
        '.docx', '.doc',
        '.xlsx', '.xls',
        '.pptx', '.ppt',
        '.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff',
        '.csv',
        '.json',
        '.html', '.htm',
        '.md', '.markdown'
    ]


def check_format_support() -> Dict:
    """Check which formats are currently supported based on installed libraries"""
    
    support = {
        'text': True,  # Always supported
        'pdf': False,
        'word': False,
        'excel': False,
        'powerpoint': False,
        'image_ocr': False,
        'csv': True,  # Always supported
        'json': True,  # Always supported
        'html': False,
        'markdown': True  # Always supported
    }
    
    # Check PDF
    try:
        import PyPDF2
        support['pdf'] = True
    except ImportError:
        try:
            import pdfplumber
            support['pdf'] = True
        except ImportError:
            pass
    
    # Check Word
    try:
        import docx
        support['word'] = True
    except ImportError:
        pass
    
    # Check Excel
    try:
        import openpyxl
        support['excel'] = True
    except ImportError:
        pass
    
    # Check PowerPoint
    try:
        from pptx import Presentation
        support['powerpoint'] = True
    except ImportError:
        pass
    
    # Check Image OCR
    try:
        from PIL import Image
        import pytesseract
        support['image_ocr'] = True
    except ImportError:
        pass
    
    # Check HTML
    try:
        from bs4 import BeautifulSoup
        support['html'] = True
    except ImportError:
        pass
    
    return support
